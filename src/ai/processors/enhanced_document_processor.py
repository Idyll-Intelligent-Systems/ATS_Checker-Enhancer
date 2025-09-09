"""
Enhanced Multi-Format Document Processor for ZeX-ATS-AI
Supports: PDF, DOCX, LaTeX, Images (JPEG/PNG), PPT, Excel, Audio (MP3), Video (MP4)
"""

import asyncio
import io
import os
import tempfile
import zipfile
from pathlib import Path
from typing import Dict, List, Optional, Union, Any
import mimetypes
import logging

# Document processing libraries
import fitz  # PyMuPDF for PDF processing
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation
import speech_recognition as sr
import cv2
from moviepy.editor import VideoFileClip
import pytesseract
from PIL import Image, ImageEnhance, ImageFilter
import numpy as np

# AI and text processing
import openai
from transformers import pipeline
import whisper

# Async file handling
import aiofiles
from fastapi import UploadFile

from src.core.config import settings
from src.utils.helpers import generate_unique_filename
from src.ai.models.openai_client import OpenAIClient


logger = logging.getLogger(__name__)


class EnhancedDocumentProcessor:
    """Enhanced document processor supporting multiple file formats."""
    
    SUPPORTED_FORMATS = {
        # Documents
        'pdf': ['application/pdf'],
        'docx': [
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'application/msword'
        ],
        'tex': ['text/x-tex', 'application/x-tex'],
        'latex': ['text/x-latex', 'application/x-latex'],
        
        # Presentations
        'pptx': [
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'application/vnd.ms-powerpoint'
        ],
        
        # Spreadsheets
        'xlsx': [
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'application/vnd.ms-excel'
        ],
        
        # Images
        'jpeg': ['image/jpeg'],
        'jpg': ['image/jpeg'],
        'png': ['image/png'],
        'tiff': ['image/tiff'],
        'bmp': ['image/bmp'],
        
        # Audio
        'mp3': ['audio/mpeg', 'audio/mp3'],
        'wav': ['audio/wav'],
        'm4a': ['audio/m4a'],
        
        # Video
        'mp4': ['video/mp4'],
        'avi': ['video/avi'],
        'mov': ['video/mov', 'video/quicktime'],
        'mkv': ['video/mkv']
    }
    
    def __init__(self):
        """Initialize the enhanced document processor."""
        self.openai_client = OpenAIClient()
        self.whisper_model = None
        self.ocr_languages = ['eng', 'spa', 'fra', 'deu', 'ita', 'por']  # Multi-language OCR
        
        # Initialize Whisper for audio transcription
        try:
            self.whisper_model = whisper.load_model("base")
            logger.info("Whisper model loaded for audio transcription")
        except Exception as e:
            logger.warning(f"Failed to load Whisper model: {e}")
        
        # Initialize speech recognition
        self.recognizer = sr.Recognizer()
        
        # Initialize AI pipeline for text analysis
        self.text_analyzer = pipeline(
            "text-classification",
            model="microsoft/DialoGPT-medium",
            return_all_scores=True
        )
    
    async def process_document(
        self, 
        file: UploadFile, 
        analysis_type: str = "resume"
    ) -> Dict[str, Any]:
        """
        Process a document of any supported format.
        
        Args:
            file: Uploaded file
            analysis_type: Type of analysis (resume, job_description, portfolio)
            
        Returns:
            Processed document data with extracted content and analysis
        """
        try:
            # Validate file format
            file_info = await self._validate_file(file)
            
            # Save file temporarily
            temp_path = await self._save_temp_file(file)
            
            try:
                # Process based on file type
                processor_method = getattr(
                    self, 
                    f"_process_{file_info['format']}", 
                    self._process_unknown
                )
                
                content_data = await processor_method(temp_path, file_info)
                
                # Enhance with AI analysis
                enhanced_data = await self._enhance_with_ai(content_data, analysis_type)
                
                # Generate insights
                insights = await self._generate_insights(enhanced_data, file_info['format'])
                
                return {
                    "success": True,
                    "file_info": file_info,
                    "content": enhanced_data,
                    "insights": insights,
                    "processing_time": content_data.get("processing_time", 0)
                }
                
            finally:
                # Clean up temporary file
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
                    
        except Exception as e:
            logger.error(f"Error processing document: {str(e)}")
            return {
                "success": False,
                "error": str(e),
                "file_info": {"filename": file.filename}
            }
    
    async def _validate_file(self, file: UploadFile) -> Dict[str, Any]:
        """Validate file format and extract basic info."""
        filename = file.filename or "unknown"
        file_extension = Path(filename).suffix.lower().lstrip('.')
        
        # Check if format is supported
        if file_extension not in self.SUPPORTED_FORMATS:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        # Validate MIME type
        expected_mimes = self.SUPPORTED_FORMATS[file_extension]
        if file.content_type and file.content_type not in expected_mimes:
            logger.warning(f"MIME type mismatch: {file.content_type} vs {expected_mimes}")
        
        # Get file size
        file.file.seek(0, 2)  # Seek to end
        file_size = file.file.tell()
        file.file.seek(0)  # Reset to beginning
        
        return {
            "filename": filename,
            "format": file_extension,
            "mime_type": file.content_type,
            "size": file_size,
            "supported": True
        }
    
    async def _save_temp_file(self, file: UploadFile) -> str:
        """Save uploaded file to temporary location."""
        suffix = Path(file.filename or "temp").suffix
        temp_fd, temp_path = tempfile.mkstemp(suffix=suffix)
        
        try:
            with os.fdopen(temp_fd, 'wb') as tmp_file:
                content = await file.read()
                tmp_file.write(content)
            
            # Reset file pointer for potential re-reading
            await file.seek(0)
            
            return temp_path
        except Exception as e:
            os.close(temp_fd)
            if os.path.exists(temp_path):
                os.unlink(temp_path)
            raise e
    
    # PDF Processing
    async def _process_pdf(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process PDF documents."""
        content_data = {
            "text": "",
            "pages": [],
            "images": [],
            "metadata": {},
            "structure": {}
        }
        
        try:
            doc = fitz.open(file_path)
            content_data["metadata"] = doc.metadata
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # Extract text
                page_text = page.get_text()
                content_data["text"] += page_text + "\n"
                
                # Extract page structure
                blocks = page.get_text("dict")
                page_data = {
                    "page_number": page_num + 1,
                    "text": page_text,
                    "blocks": len(blocks.get("blocks", [])),
                    "fonts": self._extract_fonts(blocks)
                }
                content_data["pages"].append(page_data)
                
                # Extract images from page
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    image_data = await self._extract_pdf_image(doc, img, page_num, img_index)
                    if image_data:
                        content_data["images"].append(image_data)
            
            doc.close()
            
        except Exception as e:
            logger.error(f"Error processing PDF: {e}")
            content_data["error"] = str(e)
        
        return content_data
    
    # DOCX Processing
    async def _process_docx(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process Microsoft Word documents."""
        content_data = {
            "text": "",
            "paragraphs": [],
            "tables": [],
            "images": [],
            "styles": {},
            "metadata": {}
        }
        
        try:
            doc = DocxDocument(file_path)
            
            # Extract core properties
            content_data["metadata"] = {
                "title": doc.core_properties.title,
                "author": doc.core_properties.author,
                "subject": doc.core_properties.subject,
                "created": str(doc.core_properties.created),
                "modified": str(doc.core_properties.modified)
            }
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    para_data = {
                        "text": para.text,
                        "style": para.style.name if para.style else "Normal",
                        "alignment": str(para.alignment) if para.alignment else None
                    }
                    content_data["paragraphs"].append(para_data)
                    content_data["text"] += para.text + "\n"
            
            # Extract tables
            for table_index, table in enumerate(doc.tables):
                table_data = {
                    "table_index": table_index,
                    "rows": len(table.rows),
                    "columns": len(table.columns),
                    "data": []
                }
                
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data["data"].append(row_data)
                
                content_data["tables"].append(table_data)
            
            # Extract embedded images
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_data = await self._process_embedded_image(rel)
                    if image_data:
                        content_data["images"].append(image_data)
                        
        except Exception as e:
            logger.error(f"Error processing DOCX: {e}")
            content_data["error"] = str(e)
        
        return content_data
    
    # LaTeX Processing
    async def _process_tex(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process LaTeX documents."""
        return await self._process_latex(file_path, file_info)
    
    async def _process_latex(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process LaTeX documents."""
        content_data = {
            "raw_text": "",
            "cleaned_text": "",
            "sections": [],
            "commands": {},
            "packages": [],
            "structure": {}
        }
        
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                latex_content = await f.read()
            
            content_data["raw_text"] = latex_content
            
            # Parse LaTeX structure
            structure = await self._parse_latex_structure(latex_content)
            content_data.update(structure)
            
            # Clean LaTeX commands to extract plain text
            cleaned_text = await self._clean_latex_text(latex_content)
            content_data["cleaned_text"] = cleaned_text
            
        except Exception as e:
            logger.error(f"Error processing LaTeX: {e}")
            content_data["error"] = str(e)
        
        return content_data
    
    # PowerPoint Processing
    async def _process_pptx(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process PowerPoint presentations."""
        content_data = {
            "text": "",
            "slides": [],
            "images": [],
            "metadata": {},
            "structure": {}
        }
        
        try:
            prs = Presentation(file_path)
            
            # Extract presentation metadata
            content_data["metadata"] = {
                "title": prs.core_properties.title,
                "author": prs.core_properties.author,
                "subject": prs.core_properties.subject,
                "created": str(prs.core_properties.created),
                "slide_count": len(prs.slides)
            }
            
            # Process each slide
            for slide_index, slide in enumerate(prs.slides):
                slide_data = {
                    "slide_number": slide_index + 1,
                    "title": "",
                    "content": [],
                    "images": [],
                    "layout": slide.slide_layout.name
                }
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if shape == slide.shapes.title:
                            slide_data["title"] = shape.text
                        else:
                            slide_data["content"].append(shape.text)
                        
                        content_data["text"] += shape.text + "\n"
                    
                    # Extract images
                    if hasattr(shape, "image"):
                        image_data = await self._process_ppt_image(shape.image, slide_index)
                        if image_data:
                            slide_data["images"].append(image_data)
                            content_data["images"].append(image_data)
                
                content_data["slides"].append(slide_data)
                
        except Exception as e:
            logger.error(f"Error processing PPTX: {e}")
            content_data["error"] = str(e)
        
        return content_data
    
    # Excel Processing
    async def _process_xlsx(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process Excel spreadsheets."""
        content_data = {
            "text": "",
            "worksheets": [],
            "metadata": {},
            "summary": {}
        }
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            
            # Extract workbook metadata
            content_data["metadata"] = {
                "title": workbook.properties.title,
                "creator": workbook.properties.creator,
                "created": str(workbook.properties.created),
                "modified": str(workbook.properties.modified),
                "worksheet_count": len(workbook.worksheets)
            }
            
            # Process each worksheet
            for ws in workbook.worksheets:
                worksheet_data = {
                    "name": ws.title,
                    "max_row": ws.max_row,
                    "max_column": ws.max_column,
                    "data": [],
                    "summary": {}
                }
                
                # Extract data (limit to reasonable size)
                max_rows = min(ws.max_row, 1000)  # Limit to prevent memory issues
                max_cols = min(ws.max_column, 50)
                
                for row in ws.iter_rows(min_row=1, max_row=max_rows, max_col=max_cols, values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    worksheet_data["data"].append(row_data)
                    
                    # Add non-empty cells to text content
                    for cell in row_data:
                        if cell and cell != "None":
                            content_data["text"] += cell + " "
                
                # Generate worksheet summary
                worksheet_data["summary"] = await self._analyze_excel_worksheet(worksheet_data)
                content_data["worksheets"].append(worksheet_data)
            
            workbook.close()
            
        except Exception as e:
            logger.error(f"Error processing XLSX: {e}")
            content_data["error"] = str(e)
        
        return content_data
    
    # Image Processing (JPEG, PNG, etc.)
    async def _process_jpeg(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process JPEG images."""
        return await self._process_image(file_path, file_info)
    
    async def _process_jpg(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process JPG images."""
        return await self._process_image(file_path, file_info)
    
    async def _process_png(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process PNG images."""
        return await self._process_image(file_path, file_info)
    
    async def _process_image(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process image files using OCR."""
        content_data = {
            "text": "",
            "ocr_confidence": 0,
            "image_info": {},
            "detected_elements": [],
            "enhanced_versions": {}
        }
        
        try:
            # Load image
            image = Image.open(file_path)
            
            # Get image info
            content_data["image_info"] = {
                "format": image.format,
                "mode": image.mode,
                "size": image.size,
                "has_transparency": image.mode in ("RGBA", "LA") or "transparency" in image.info
            }
            
            # Enhance image for better OCR
            enhanced_image = await self._enhance_image_for_ocr(image)
            
            # Perform OCR
            ocr_results = await self._perform_ocr(enhanced_image)
            content_data.update(ocr_results)
            
            # Detect document structure (resume sections, etc.)
            structure = await self._detect_document_structure(image, content_data["text"])
            content_data["detected_elements"] = structure
            
        except Exception as e:
            logger.error(f"Error processing image: {e}")
            content_data["error"] = str(e)
        
        return content_data
    
    # Audio Processing (MP3, etc.)
    async def _process_mp3(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process MP3 audio files."""
        return await self._process_audio(file_path, file_info)
    
    async def _process_wav(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process WAV audio files."""
        return await self._process_audio(file_path, file_info)
    
    async def _process_audio(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process audio files using speech recognition."""
        content_data = {
            "text": "",
            "transcription": "",
            "audio_info": {},
            "confidence": 0,
            "language": "unknown",
            "segments": []
        }
        
        try:
            # Get audio info
            if file_path.endswith(('.mp3', '.wav', '.m4a')):
                # Use Whisper for transcription if available
                if self.whisper_model:
                    result = self.whisper_model.transcribe(file_path)
                    content_data["transcription"] = result["text"]
                    content_data["text"] = result["text"]
                    content_data["language"] = result.get("language", "unknown")
                    content_data["segments"] = result.get("segments", [])
                else:
                    # Fallback to speech_recognition
                    with sr.AudioFile(file_path) as source:
                        audio = self.recognizer.record(source)
                        try:
                            text = self.recognizer.recognize_google(audio)
                            content_data["transcription"] = text
                            content_data["text"] = text
                        except sr.UnknownValueError:
                            content_data["text"] = ""
                            content_data["error"] = "Could not understand audio"
                        except sr.RequestError as e:
                            content_data["error"] = f"Could not request results: {e}"
            
        except Exception as e:
            logger.error(f"Error processing audio: {e}")
            content_data["error"] = str(e)
        
        return content_data
    
    # Video Processing (MP4, etc.)
    async def _process_mp4(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process MP4 video files."""
        return await self._process_video(file_path, file_info)
    
    async def _process_avi(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process AVI video files."""
        return await self._process_video(file_path, file_info)
    
    async def _process_mov(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process MOV video files."""
        return await self._process_video(file_path, file_info)
    
    async def _process_video(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Process video files - extract audio and analyze frames."""
        content_data = {
            "text": "",
            "audio_transcription": "",
            "video_info": {},
            "key_frames": [],
            "ocr_results": [],
            "summary": ""
        }
        
        try:
            # Load video
            video = VideoFileClip(file_path)
            
            content_data["video_info"] = {
                "duration": video.duration,
                "fps": video.fps,
                "size": video.size,
                "has_audio": video.audio is not None
            }
            
            # Extract and transcribe audio if present
            if video.audio:
                audio_path = tempfile.mktemp(suffix=".wav")
                try:
                    video.audio.write_audiofile(audio_path, verbose=False, logger=None)
                    audio_result = await self._process_audio(audio_path, {"format": "wav"})
                    content_data["audio_transcription"] = audio_result.get("text", "")
                    content_data["text"] += audio_result.get("text", "")
                finally:
                    if os.path.exists(audio_path):
                        os.unlink(audio_path)
            
            # Extract key frames and perform OCR
            duration = video.duration
            frame_times = [duration * i / 10 for i in range(0, 11)]  # 11 frames
            
            for i, time in enumerate(frame_times):
                try:
                    frame = video.get_frame(time)
                    frame_image = Image.fromarray(frame.astype('uint8'), 'RGB')
                    
                    # Perform OCR on frame
                    frame_ocr = await self._perform_ocr(frame_image)
                    if frame_ocr.get("text", "").strip():
                        content_data["ocr_results"].append({
                            "frame": i,
                            "time": time,
                            "text": frame_ocr["text"],
                            "confidence": frame_ocr.get("ocr_confidence", 0)
                        })
                        content_data["text"] += frame_ocr["text"] + " "
                
                except Exception as e:
                    logger.warning(f"Error processing frame {i}: {e}")
            
            video.close()
            
        except Exception as e:
            logger.error(f"Error processing video: {e}")
            content_data["error"] = str(e)
        
        return content_data
    
    # Helper Methods
    async def _enhance_image_for_ocr(self, image: Image.Image) -> Image.Image:
        """Enhance image quality for better OCR results."""
        try:
            # Convert to grayscale if not already
            if image.mode != 'L':
                image = image.convert('L')
            
            # Enhance contrast
            enhancer = ImageEnhance.Contrast(image)
            image = enhancer.enhance(2.0)
            
            # Enhance sharpness
            enhancer = ImageEnhance.Sharpness(image)
            image = enhancer.enhance(2.0)
            
            # Apply slight blur to reduce noise
            image = image.filter(ImageFilter.MedianFilter(size=3))
            
            return image
            
        except Exception as e:
            logger.warning(f"Error enhancing image: {e}")
            return image
    
    async def _perform_ocr(self, image: Image.Image) -> Dict[str, Any]:
        """Perform OCR on image using Tesseract."""
        ocr_results = {
            "text": "",
            "ocr_confidence": 0,
            "detected_languages": []
        }
        
        try:
            # Configure Tesseract
            custom_config = r'--oem 3 --psm 6 -l eng+spa+fra+deu+ita+por'
            
            # Get detailed OCR data
            ocr_data = pytesseract.image_to_data(
                image, 
                config=custom_config, 
                output_type=pytesseract.Output.DICT
            )
            
            # Extract text and confidence
            words = []
            confidences = []
            
            for i in range(len(ocr_data['text'])):
                word = ocr_data['text'][i].strip()
                conf = int(ocr_data['conf'][i])
                
                if word and conf > 30:  # Filter low confidence words
                    words.append(word)
                    confidences.append(conf)
            
            ocr_results["text"] = " ".join(words)
            ocr_results["ocr_confidence"] = sum(confidences) / len(confidences) if confidences else 0
            
        except Exception as e:
            logger.error(f"OCR error: {e}")
            # Fallback to simple OCR
            try:
                ocr_results["text"] = pytesseract.image_to_string(image)
                ocr_results["ocr_confidence"] = 50  # Default confidence
            except Exception as fallback_error:
                logger.error(f"Fallback OCR error: {fallback_error}")
        
        return ocr_results
    
    async def _detect_document_structure(self, image: Image.Image, text: str) -> List[Dict]:
        """Detect document structure elements (headers, sections, etc.)."""
        elements = []
        
        try:
            # Use OpenCV for structure detection
            img_array = np.array(image)
            
            # Convert to grayscale if needed
            if len(img_array.shape) == 3:
                gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
            else:
                gray = img_array
            
            # Detect horizontal lines (section separators)
            horizontal_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (25, 1))
            horizontal_lines = cv2.morphologyEx(gray, cv2.MORPH_OPEN, horizontal_kernel)
            
            # Find contours of horizontal lines
            contours, _ = cv2.findContours(horizontal_lines, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            for contour in contours:
                x, y, w, h = cv2.boundingRect(contour)
                if w > image.width * 0.3:  # Line spans at least 30% of width
                    elements.append({
                        "type": "horizontal_line",
                        "position": {"x": x, "y": y, "width": w, "height": h}
                    })
            
            # Analyze text structure
            if text:
                lines = text.split('\n')
                for i, line in enumerate(lines):
                    line = line.strip()
                    if line:
                        # Detect potential headers (short lines, all caps, etc.)
                        if len(line) < 50 and (line.isupper() or ':' in line):
                            elements.append({
                                "type": "header",
                                "text": line,
                                "line_number": i
                            })
                        
                        # Detect email patterns
                        if '@' in line and '.' in line:
                            elements.append({
                                "type": "email",
                                "text": line,
                                "line_number": i
                            })
                        
                        # Detect phone patterns
                        import re
                        phone_pattern = r'(\+\d{1,3}[-.\s]?)?\(?\d{3}\)?[-.\s]?\d{3}[-.\s]?\d{4}'
                        if re.search(phone_pattern, line):
                            elements.append({
                                "type": "phone",
                                "text": line,
                                "line_number": i
                            })
        
        except Exception as e:
            logger.warning(f"Error detecting document structure: {e}")
        
        return elements
    
    async def _parse_latex_structure(self, latex_content: str) -> Dict[str, Any]:
        """Parse LaTeX document structure."""
        structure = {
            "document_class": "",
            "packages": [],
            "sections": [],
            "commands": {},
            "bibliography": False
        }
        
        try:
            import re
            
            # Extract document class
            doc_class_match = re.search(r'\\documentclass(?:\[.*?\])?\{(.*?)\}', latex_content)
            if doc_class_match:
                structure["document_class"] = doc_class_match.group(1)
            
            # Extract packages
            package_matches = re.findall(r'\\usepackage(?:\[.*?\])?\{(.*?)\}', latex_content)
            structure["packages"] = package_matches
            
            # Extract sections
            section_patterns = [
                (r'\\section\{(.*?)\}', 'section'),
                (r'\\subsection\{(.*?)\}', 'subsection'),
                (r'\\subsubsection\{(.*?)\}', 'subsubsection'),
                (r'\\chapter\{(.*?)\}', 'chapter')
            ]
            
            for pattern, section_type in section_patterns:
                matches = re.finditer(pattern, latex_content)
                for match in matches:
                    structure["sections"].append({
                        "type": section_type,
                        "title": match.group(1),
                        "position": match.start()
                    })
            
            # Check for bibliography
            if '\\bibliography' in latex_content or '\\bibitem' in latex_content:
                structure["bibliography"] = True
            
        except Exception as e:
            logger.warning(f"Error parsing LaTeX structure: {e}")
        
        return structure
    
    async def _clean_latex_text(self, latex_content: str) -> str:
        """Remove LaTeX commands and extract plain text."""
        try:
            import re
            
            # Remove comments
            text = re.sub(r'%.*', '', latex_content)
            
            # Remove common commands but keep their content
            commands_with_content = [
                r'\\textbf\{(.*?)\}',
                r'\\textit\{(.*?)\}',
                r'\\emph\{(.*?)\}',
                r'\\section\{(.*?)\}',
                r'\\subsection\{(.*?)\}',
                r'\\title\{(.*?)\}',
                r'\\author\{(.*?)\}',
            ]
            
            for pattern in commands_with_content:
                text = re.sub(pattern, r'\1', text)
            
            # Remove commands without content
            text = re.sub(r'\\[a-zA-Z]+(\[.*?\])?(\{.*?\})?', '', text)
            
            # Remove special characters
            text = re.sub(r'[{}\\$^_~]', '', text)
            
            # Clean up whitespace
            text = re.sub(r'\s+', ' ', text)
            text = text.strip()
            
            return text
            
        except Exception as e:
            logger.warning(f"Error cleaning LaTeX text: {e}")
            return latex_content
    
    async def _enhance_with_ai(self, content_data: Dict, analysis_type: str) -> Dict[str, Any]:
        """Enhance extracted content with AI analysis."""
        enhanced_data = content_data.copy()
        
        try:
            text = content_data.get("text", "")
            if not text.strip():
                return enhanced_data
            
            # Generate AI insights based on analysis type
            if analysis_type == "resume":
                ai_analysis = await self._analyze_resume_content(text)
            elif analysis_type == "job_description":
                ai_analysis = await self._analyze_job_description(text)
            elif analysis_type == "portfolio":
                ai_analysis = await self._analyze_portfolio_content(text)
            else:
                ai_analysis = await self._general_text_analysis(text)
            
            enhanced_data["ai_analysis"] = ai_analysis
            
        except Exception as e:
            logger.error(f"Error enhancing with AI: {e}")
            enhanced_data["ai_error"] = str(e)
        
        return enhanced_data
    
    async def _analyze_resume_content(self, text: str) -> Dict[str, Any]:
        """AI analysis specifically for resume content."""
        try:
            prompt = f"""
            Analyze this resume content and provide structured insights:
            
            Resume Text: {text[:3000]}...
            
            Please provide:
            1. Key skills mentioned
            2. Experience level assessment
            3. Industries/domains identified
            4. Strengths and areas for improvement
            5. ATS optimization suggestions
            6. Missing critical elements
            
            Respond in JSON format with clear categories.
            """
            
            response = await self.openai_client.generate_completion(prompt)
            
            # Parse AI response
            import json
            try:
                ai_insights = json.loads(response)
            except:
                # Fallback if JSON parsing fails
                ai_insights = {
                    "analysis": response,
                    "skills": [],
                    "experience_level": "unknown",
                    "recommendations": []
                }
            
            return ai_insights
            
        except Exception as e:
            logger.error(f"Error in resume AI analysis: {e}")
            return {"error": str(e)}
    
    async def _generate_insights(self, enhanced_data: Dict, format_type: str) -> Dict[str, Any]:
        """Generate format-specific insights."""
        insights = {
            "format_specific": {},
            "content_quality": {},
            "recommendations": [],
            "ats_score": 0
        }
        
        try:
            text = enhanced_data.get("text", "")
            
            # Format-specific insights
            if format_type in ["pdf", "docx"]:
                insights["format_specific"] = {
                    "readability": "good" if len(text) > 100 else "poor",
                    "structure": "organized" if enhanced_data.get("pages") or enhanced_data.get("paragraphs") else "unstructured",
                    "formatting": "professional"
                }
            
            elif format_type in ["jpeg", "png", "jpg"]:
                insights["format_specific"] = {
                    "ocr_quality": enhanced_data.get("ocr_confidence", 0),
                    "image_quality": "good" if enhanced_data.get("ocr_confidence", 0) > 70 else "needs_improvement",
                    "text_extraction": "successful" if text else "failed"
                }
            
            elif format_type in ["mp3", "wav", "m4a"]:
                insights["format_specific"] = {
                    "transcription_quality": "good" if text else "poor",
                    "audio_clarity": "clear",
                    "language_detected": enhanced_data.get("language", "unknown")
                }
            
            elif format_type in ["mp4", "avi", "mov"]:
                insights["format_specific"] = {
                    "content_types": ["audio", "visual"] if enhanced_data.get("audio_transcription") and enhanced_data.get("ocr_results") else ["visual"],
                    "extraction_success": "partial" if text else "minimal"
                }
            
            # Calculate basic ATS score
            insights["ats_score"] = await self._calculate_basic_ats_score(enhanced_data)
            
            # Generate recommendations
            insights["recommendations"] = await self._generate_recommendations(enhanced_data, format_type)
            
        except Exception as e:
            logger.error(f"Error generating insights: {e}")
            insights["error"] = str(e)
        
        return insights
    
    async def _calculate_basic_ats_score(self, data: Dict) -> int:
        """Calculate a basic ATS compatibility score."""
        score = 0
        text = data.get("text", "")
        
        try:
            # Text length (20 points)
            if len(text) > 500:
                score += 20
            elif len(text) > 200:
                score += 10
            
            # Contact information (20 points)
            if '@' in text and '.' in text:  # Email
                score += 10
            
            import re
            if re.search(r'\d{3}[-.\s]?\d{3}[-.\s]?\d{4}', text):  # Phone
                score += 10
            
            # Professional keywords (30 points)
            keywords = ['experience', 'skills', 'education', 'work', 'professional', 'project', 'achievement']
            found_keywords = sum(1 for keyword in keywords if keyword.lower() in text.lower())
            score += min(found_keywords * 5, 30)
            
            # Structure (30 points)
            if data.get("paragraphs") or data.get("sections"):
                score += 20
            if data.get("tables") or data.get("pages"):
                score += 10
            
            return min(score, 100)
            
        except Exception as e:
            logger.warning(f"Error calculating ATS score: {e}")
            return 50  # Default score
    
    async def _generate_recommendations(self, data: Dict, format_type: str) -> List[str]:
        """Generate improvement recommendations."""
        recommendations = []
        text = data.get("text", "")
        
        try:
            # Format-specific recommendations
            if format_type in ["jpeg", "png", "jpg"] and data.get("ocr_confidence", 0) < 70:
                recommendations.append("Consider using a text-based format (PDF, DOCX) for better readability")
                recommendations.append("If using images, ensure high resolution and clear text")
            
            if format_type in ["mp3", "wav"] and not text:
                recommendations.append("Audio transcription failed - consider providing a text version")
            
            if format_type == "latex":
                recommendations.append("LaTeX detected - consider converting to PDF for broader compatibility")
            
            # Content recommendations
            if len(text) < 200:
                recommendations.append("Document appears to be very short - consider adding more detail")
            
            if '@' not in text:
                recommendations.append("Add contact email address for better ATS parsing")
            
            if not any(word in text.lower() for word in ['experience', 'work', 'job']):
                recommendations.append("Include work experience section")
            
            if not any(word in text.lower() for word in ['skill', 'ability', 'proficient']):
                recommendations.append("Add skills section with relevant keywords")
            
        except Exception as e:
            logger.warning(f"Error generating recommendations: {e}")
        
        return recommendations[:5]  # Limit to top 5 recommendations
    
    async def _process_unknown(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """Handle unknown file formats."""
        return {
            "text": "",
            "error": f"Unsupported file format: {file_info.get('format', 'unknown')}",
            "supported_formats": list(self.SUPPORTED_FORMATS.keys())
        }
    
    # Additional helper methods for specific file types
    async def _extract_fonts(self, blocks: Dict) -> List[str]:
        """Extract font information from PDF blocks."""
        fonts = set()
        try:
            for block in blocks.get("blocks", []):
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        if "font" in span:
                            fonts.add(span["font"])
        except Exception as e:
            logger.warning(f"Error extracting fonts: {e}")
        return list(fonts)
    
    async def _extract_pdf_image(self, doc, img, page_num: int, img_index: int) -> Optional[Dict]:
        """Extract image from PDF document."""
        try:
            xref = img[0]
            pix = fitz.Pixmap(doc, xref)
            
            if pix.n - pix.alpha < 4:  # Valid image
                img_data = pix.tobytes("png")
                return {
                    "page": page_num + 1,
                    "index": img_index,
                    "size": (pix.width, pix.height),
                    "colorspace": pix.colorspace.name if pix.colorspace else "unknown"
                }
            
            pix = None
            return None
            
        except Exception as e:
            logger.warning(f"Error extracting PDF image: {e}")
            return None
